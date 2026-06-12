import os
import re

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree


class PresentationService:

    def __init__(self):
        os.makedirs("outputs", exist_ok=True)
        os.makedirs("outputs/charts", exist_ok=True)

    def create_report(
        self,
        file_name,
        profile,
        quality_report,
        analysis_result,
        insights,
        recommendations,
        chart_items=None,
        output_path="outputs/AI_Report.pptx"
    ):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        base_name = os.path.splitext(file_name)[0].replace("_", " ").replace("-", " ").title()
        self._current_base_name = base_name

        # Determine dynamic titles based on slide content or uploaded file
        insights_default = f"{base_name} Insights" if base_name else "Business Insights"
        recs_default = f"{base_name} Recommendations" if base_name else "Recommendations"

        insights_title = self._determine_slide_title(insights, insights_default)
        recs_title = self._determine_slide_title(recommendations, recs_default)

        self._add_title_slide(prs, file_name)
        self._add_dataset_overview(prs, profile)
        self._add_data_quality_slide(prs, quality_report)

        # Only add slides if they contain valid generated content (skipping placeholders/errors)
        if self._is_valid_content(analysis_result):
            self._add_analysis_slide(prs, analysis_result)

        if chart_items:
            self._add_visualization_slides(prs, chart_items)

        if self._is_valid_content(insights):
            final_ins_title = insights_title if insights_title.startswith("💡") else f"💡 {insights_title}"
            self._add_text_slides(prs, final_ins_title, insights)

        if self._is_valid_content(recommendations):
            final_rec_title = recs_title if recs_title.startswith("🎯") else f"🎯 {recs_title}"
            self._add_text_slides(prs, final_rec_title, recommendations)

        self._add_conclusion_slide(prs)

        self._remove_empty_slides(prs)
        prs.save(output_path)
        return output_path

    def _blank_slide(self, prs, title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 250, 252) # slate-50 background

        # Add a subtle header background band
        header_band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(1.1),
        )
        header_band.fill.solid()
        header_band.fill.fore_color.rgb = RGBColor(241, 245, 249) # slate-100 header band
        header_band.line.fill.background()

        # Add a colored left vertical accent line
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.4),
            Inches(0.32),
            Inches(0.08),
            Inches(0.46),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(13, 148, 136) # Teal accent color
        accent_bar.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.65),
            Inches(0.22),
            Inches(12),
            Inches(0.65),
        )
        frame = title_box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        # Clean title of any markdown artifacts
        clean_title = str(title)
        clean_title = re.sub(r"\*\*(.*?)\*\*", r"\1", clean_title)
        clean_title = re.sub(r"^#+\s*", "", clean_title).strip(":- ")
        p.text = clean_title
        p.font.name = "Times New Roman"
        p.font.size = Pt(28 if len(clean_title) <= 42 else 22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 41, 59) # Slate 800 title text
        return slide

    def _get_emoji_for_text(self, text):
        text_lower = text.lower()
        if "insight" in text_lower or "found" in text_lower:
            return "💡"
        if "recommend" in text_lower or "action" in text_lower or "should" in text_lower:
            return "🎯"
        if "trend" in text_lower or "growth" in text_lower or "increase" in text_lower or "rise" in text_lower:
            return "📈"
        if "decrease" in text_lower or "decline" in text_lower or "drop" in text_lower or "loss" in text_lower:
            return "📉"
        if "missing" in text_lower or "null" in text_lower or "empty" in text_lower or "blank" in text_lower or "nan" in text_lower:
            return "🔍"
        if "duplicate" in text_lower or "repeat" in text_lower or "redundant" in text_lower:
            return "🔄"
        if "error" in text_lower or "fail" in text_lower or "wrong" in text_lower or "invalid" in text_lower or "issue" in text_lower or "warning" in text_lower or "caution" in text_lower:
            return "⚠️"
        if "number" in text_lower or "count" in text_lower or "sum" in text_lower or "total" in text_lower or "average" in text_lower or "mean" in text_lower or "median" in text_lower or "percent" in text_lower or "%" in text_lower:
            return "📊"
        if "check" in text_lower or "done" in text_lower or "success" in text_lower or "verify" in text_lower or "complete" in text_lower or "resolve" in text_lower:
            return "✅"
        if "important" in text_lower or "key" in text_lower or "critical" in text_lower or "significant" in text_lower or "essential" in text_lower:
            return "🔑"
        if "customer" in text_lower or "user" in text_lower or "segment" in text_lower or "people" in text_lower or "demographic" in text_lower:
            return "👥"
        if "cost" in text_lower or "price" in text_lower or "expense" in text_lower or "spend" in text_lower or "revenue" in text_lower or "profit" in text_lower or "sales" in text_lower or "financial" in text_lower or "dollar" in text_lower or "amount" in text_lower:
            return "💰"
        if "time" in text_lower or "date" in text_lower or "month" in text_lower or "year" in text_lower or "weekly" in text_lower or "daily" in text_lower or "hourly" in text_lower:
            return "📅"
        if "speed" in text_lower or "fast" in text_lower or "slow" in text_lower or "performance" in text_lower or "optimize" in text_lower or "quick" in text_lower or "accelerate" in text_lower:
            return "⚡"
        if "compare" in text_lower or "versus" in text_lower or "vs" in text_lower or "difference" in text_lower or "diff" in text_lower:
            return "⚖️"
        if "future" in text_lower or "predict" in text_lower or "forecast" in text_lower or "projection" in text_lower:
            return "🔮"
        return "•"

    def create_report_from_state(self, presentation_state, output_path="outputs/Editable_AI_Report.pptx"):
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs = Presentation()
        
        for slide_data in presentation_state.get("slides", []):
            title = slide_data.get("title", "Slide")
            
            text_lines = []
            for elem in slide_data.get("elements", []):
                if elem.get("type") in ["textbox", "text"]:
                    text_lines.append(elem.get("text", ""))
                
            combined_text = "\n".join(text_lines)
            
            if slide_data.get("layout") == "title_slide":
                # Hack to pass title as file_name to title slide generator
                self._add_title_slide(prs, title)
            else:
                self._add_text_slides(prs, title, combined_text)
                
        self._remove_empty_slides(prs)
        prs.save(output_path)
        return output_path

    def _add_kpi_card(self, slide, left, top, width, height, title, value, emoji=""):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(226, 232, 240)
        card.line.width = Pt(1.5)
        
        value_box = slide.shapes.add_textbox(
            left + Inches(0.15),
            top + Inches(0.15),
            width - Inches(0.3),
            Inches(0.9)
        )
        tf_val = value_box.text_frame
        tf_val.word_wrap = True
        tf_val.clear()
        p_val = tf_val.paragraphs[0]
        p_val.text = f"{emoji} {value}" if emoji else str(value)
        p_val.font.name = "Times New Roman"
        p_val.font.size = Pt(34)
        p_val.font.bold = True
        p_val.font.color.rgb = RGBColor(13, 148, 136) # Teal value accent
        
        title_box = slide.shapes.add_textbox(
            left + Inches(0.15),
            top + Inches(1.05),
            width - Inches(0.3),
            Inches(0.6)
        )
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.clear()
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Times New Roman"
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(71, 85, 105) # Slate-600

    def _get_asset_path(self, filename):
        current_dir = os.path.dirname(os.path.abspath(__file__))
        workspace_dir = os.path.dirname(current_dir)
        path = os.path.join(workspace_dir, "assets", filename)
        if os.path.exists(path):
            return path
        fallback = os.path.join("assets", filename)
        if os.path.exists(fallback):
            return fallback
        return None

    def _add_illustration_card(self, slide, left, top, width, height, asset_name):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(226, 232, 240)
        card.line.width = Pt(1.5)
        
        path = self._get_asset_path(asset_name)
        if path:
            # Use generous margins (1.2 inches total margin) to decrease image size
            pic_w_max = width - Inches(1.2)
            pic_h_max = height - Inches(1.2)
            
            # Since icons are square, size is the minimum of w/h capacity
            pic_size = min(pic_w_max, pic_h_max)
            
            # Center the picture perfectly inside the card shape
            pic_left = left + (width - pic_size) / 2
            pic_top = top + (height - pic_size) / 2
            
            try:
                slide.shapes.add_picture(
                    path,
                    pic_left,
                    pic_top,
                    width=pic_size,
                    height=pic_size
                )
            except Exception:
                pass

    def _clean_text(self, value):
        if value is None:
            text = ""
        elif isinstance(value, pd.DataFrame):
            text = value.head(20).to_string(index=True)
        elif isinstance(value, pd.Series):
            text = value.head(30).to_string()
        else:
            text = str(value)

        text = re.sub(r"`([^`]*)`", r"\1", text)
        text = text.replace("Ã¢â‚¬Â¢", "-")
        text = text.replace("\u00e2\u0080\u0093", "-")
        text = text.replace("\u00e2\u0080\u0094", "-")
        text = text.replace("\u2022", "-")
        text = re.sub(
            r"(Insight|Recommendation|Finding)\s+generation failed:\s+503.*",
            (
                r"\1 generation is temporarily unavailable because the AI model "
                "endpoint is under high demand. Please retry this section."
            ),
            text,
            flags=re.IGNORECASE | re.DOTALL,
        )

        lines = []
        for raw in text.splitlines():
            line = raw.strip()
            if not line:
                continue
            line = re.sub(r"^[\-\*\u2022]\s*", "- ", line)
            line = re.sub(r"^\d+[\.\)]\s*", "- ", line)
            lines.append(line)

        return lines

    def _title_from_caption(self, caption, fallback):
        cleaned = self._clean_text(caption)
        if not cleaned:
            return fallback

        title = cleaned[0].split(":", 1)[0].strip()
        return title[:74] or fallback

    def _is_valid_content(self, text):
        if text is None:
            return False
            
        # Handle pandas DataFrame / Series safely to avoid truthiness errors
        if isinstance(text, (pd.DataFrame, pd.Series)):
            return not text.empty
            
        if isinstance(text, dict):
            text = text.get("result", "")
            
        # Re-check in case the dictionary returned a DataFrame or Series
        if isinstance(text, (pd.DataFrame, pd.Series)):
            return not text.empty
            
        if text is None:
            return False

        # If it's not a string, check if it's truthy, otherwise convert to string
        if not isinstance(text, str):
            try:
                # Try simple comparison, if it fails, default to True (meaning it exists)
                if not text:
                    return False
            except Exception:
                pass
            try:
                text = str(text)
            except Exception:
                return True
            
        lower_text = text.strip().lower()
        placeholders = [
            "no ai insights generated yet",
            "no ai recommendations generated yet",
            "no ai analysis generated yet",
            "no details were generated for this section",
            "generation failed",
            "temporarily unavailable",
            "name 'len' is not defined",
            "is not defined",
            "error during analysis",
            "none",
            "null",
            "traceback",
            "exception",
            "error",
            "unterminated string literal",
            "syntaxerror",
            "nameerror",
        ]
        for placeholder in placeholders:
            if placeholder in lower_text:
                return False
                
        cleaned = self._clean_text(text)
        if not cleaned:
            return False
            
        return True

    def _determine_slide_title(self, text, default_title):
        if text is None:
            return default_title
            
        # Handle pandas DataFrame / Series safely to avoid truthiness errors
        if isinstance(text, (pd.DataFrame, pd.Series)):
            return default_title
            
        if isinstance(text, dict):
            text = text.get("result", "")
            
        # Re-check in case the dictionary returned a DataFrame or Series
        if isinstance(text, (pd.DataFrame, pd.Series)):
            return default_title
            
        if text is None:
            return default_title
            
        if not isinstance(text, str):
            return default_title
            
        # Try to find a markdown header (# Header, ## Header, ### Header)
        match = re.search(r'^\s*#+\s*(.+)$', text, re.MULTILINE)
        if match:
            candidate = match.group(1).strip()
            candidate = re.sub(r"\*\*(.*?)\*\*", r"\1", candidate)
            candidate = re.sub(r"\*(.*?)\*", r"\1", candidate)
            candidate = re.sub(r"`([^`]*)`", r"\1", candidate)
            candidate = candidate.strip(":- ")
            if candidate and len(candidate) < 60:
                return candidate
                
        # If no markdown header, look at the first line of cleaned text
        cleaned_lines = self._clean_text(text)
        if cleaned_lines:
            first_line = cleaned_lines[0]
            if not first_line.startswith("-") and len(first_line) < 60:
                candidate = first_line.split(":", 1)[0].strip()
                if candidate and len(candidate) > 3 and len(candidate) < 50:
                    return candidate
                    
        return default_title

    def _wrapped_lines(self, lines, width=88):
        # Do not physically wrap — let PPTX word_wrap handle layout.
        return lines

    def _set_para_indent(self, para_elem, left_emu, first_line_emu=0):
        """Set left_indent and first_line_indent on a CT_P (lxml element)."""
        pPr = para_elem.find(qn("a:pPr"))
        if pPr is None:
            pPr = etree.SubElement(para_elem, qn("a:pPr"))
            para_elem.insert(0, pPr)
        pPr.set("marL", str(int(left_emu)))
        pPr.set("indent", str(int(first_line_emu)))

    def _set_para_space_after(self, para_elem, pt_value):
        """Set spaceAft on a CT_P element in hundredths of a point."""
        pPr = para_elem.find(qn("a:pPr"))
        if pPr is None:
            pPr = etree.SubElement(para_elem, qn("a:pPr"))
            para_elem.insert(0, pPr)
        spcAft = pPr.find(qn("a:spcAft"))
        if spcAft is None:
            spcAft = etree.SubElement(pPr, qn("a:spcAft"))
        spcPts = spcAft.find(qn("a:spcPts"))
        if spcPts is None:
            spcPts = etree.SubElement(spcAft, qn("a:spcPts"))
        # hundredths of a point
        spcPts.set("val", str(int(pt_value * 100)))

    def _estimate_and_fit_text(self, lines, width_in, height_in, initial_font_size, space_after_pt):
        # Subtract margins and bullet hanging indents from usable width
        usable_width_in = max(1.0, width_in - 0.6)
        # Keep a 25pt safety threshold at the bottom of the card to prevent border crossings
        usable_height_pt = (height_in * 72) - 25

        font_size = initial_font_size
        while font_size >= 9.5:
            total_height_pt = 0
            # Times New Roman average character width is around 0.40 * font_size
            chars_per_line = int((usable_width_in * 72) / (0.40 * font_size))
            if chars_per_line <= 0:
                chars_per_line = 1
            for line in lines:
                clean_line = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
                clean_line = re.sub(r"^#+\s*", "", clean_line)
                length = len(clean_line)
                visual_lines = max(1, (length + chars_per_line - 1) // chars_per_line)
                para_height = visual_lines * (font_size * 1.3) + space_after_pt
                total_height_pt += para_height
            if total_height_pt <= usable_height_pt:
                return font_size, lines
            font_size -= 0.5
            
        font_size = 9.5
        chars_per_line = int((usable_width_in * 72) / (0.40 * font_size))
        if chars_per_line <= 0:
            chars_per_line = 1
        fitted_lines = []
        total_height_pt = 0
        for line in lines:
            clean_line = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
            clean_line = re.sub(r"^#+\s*", "", clean_line)
            length = len(clean_line)
            visual_lines = max(1, (length + chars_per_line - 1) // chars_per_line)
            para_height = visual_lines * (font_size * 1.3) + space_after_pt
            if total_height_pt + para_height <= usable_height_pt:
                total_height_pt += para_height
                fitted_lines.append(line)
            else:
                break
        if not fitted_lines and lines:
            fitted_lines = [lines[0]]
        return font_size, fitted_lines

    def _add_textbox_column(self, slide, lines, left, top, width, height, computed_font_size, space_after):
        # Prevent boundary crossing
        width_in = float(width) / Inches(1) if isinstance(width, (int, float)) else float(width)
        height_in = float(height) / Inches(1) if isinstance(height, (int, float)) else float(height)
        
        fitted_font_size, fitted_lines = self._estimate_and_fit_text(
            lines, width_in, height_in, computed_font_size, space_after
        )
        
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        frame.margin_left = Inches(0.08)
        frame.margin_right = Inches(0.08)
        frame.margin_top = Inches(0.05)
        frame.margin_bottom = Inches(0.05)

        for i, line in enumerate(fitted_lines):
            p_obj = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            p_elem = p_obj._p

            # Clean any remaining markdown bold markers or headers
            clean_line = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
            clean_line = re.sub(r"^#+\s*", "", clean_line)

            is_bullet = clean_line.startswith("- ")
            if is_bullet:
                bullet_content = clean_line[2:].strip()
                emoji_bullet = self._get_emoji_for_text(bullet_content)
                p_obj.text = f"{emoji_bullet}  {bullet_content}"
                self._set_para_indent(
                    p_elem,
                    left_emu=Inches(0.42),
                    first_line_emu=-Inches(0.26),
                )
            else:
                p_obj.text = clean_line
                self._set_para_indent(
                    p_elem,
                    left_emu=Inches(0.15),
                    first_line_emu=0,
                )

            p_obj.font.name = "Times New Roman"
            p_obj.font.size = Pt(fitted_font_size)
            p_obj.font.color.rgb = RGBColor(51, 65, 85) # Slate 700 text color
            self._set_para_space_after(p_elem, pt_value=space_after)

            if not is_bullet and not clean_line.startswith("  ") and len(clean_line) < 58:
                p_obj.font.bold = True
                p_obj.font.color.rgb = RGBColor(15, 23, 42) # Slate 900 bold headings

    def _add_body_lines(self, slide, lines, top=1.22, font_size=17, max_lines=14):
        # Truncate lines to prevent overflow
        lines = lines[:40]

        has_intro = False
        intro_line = None
        body_lines = lines

        if lines and not lines[0].startswith("- "):
            has_intro = True
            intro_line = lines[0]
            body_lines = lines[1:]

        # Clean all lines of asterisks (broken bold markdown)
        if has_intro:
            intro_line = intro_line.replace("*", "")
        body_lines = [b.replace("*", "") for b in body_lines]

        # Draw intro block if present
        intro_height = 0.0
        if has_intro:
            intro_height = 1.25
            card_intro = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.65),
                Inches(top + 0.08),
                Inches(12.03),
                Inches(intro_height)
            )
            card_intro.fill.solid()
            card_intro.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card_intro.line.color.rgb = RGBColor(226, 232, 240)
            card_intro.line.width = Pt(1.5)

            # Text inside intro
            tb_intro = slide.shapes.add_textbox(
                Inches(0.95),
                Inches(top + 0.18),
                Inches(11.43),
                Inches(intro_height - 0.2)
            )
            tf_intro = tb_intro.text_frame
            tf_intro.word_wrap = True
            tf_intro.clear()
            p_intro = tf_intro.paragraphs[0]
            p_intro.text = intro_line
            p_intro.font.name = "Times New Roman"
            p_intro.font.size = Pt(15.5)
            p_intro.font.color.rgb = RGBColor(71, 85, 105) # Slate 600
            p_intro.font.italic = True
            
            # Update top for columns
            top = top + intro_height + 0.15

        # Now handle body_lines layout
        avail_height = 5.45 - (intro_height + 0.15 if has_intro else 0)

        # Extract slide title to select topic-specific icon on single-column layouts
        slide_title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                txt = shape.text_frame.text.lower()
                if "overview" not in txt and "assessment" not in txt:
                    slide_title = txt
                    break

        icon_name = "icon_findings.png"
        if "insight" in slide_title:
            icon_name = "icon_trends.png"
        elif "recommend" in slide_title or "opportunity" in slide_title or "strategy" in slide_title:
            icon_name = "icon_opportunities.png"
        elif "conclusion" in slide_title:
            icon_name = "icon_conclusion.png"

        # Decide column count
        has_subheading = any(b.endswith(":") and len(b) < 45 for b in body_lines)
        if len(body_lines) <= 5 and not has_subheading:
            # Single column layout: Left card for text, right card for graphic illustration
            computed_font_size = 15.5 if has_intro else 17.5
            space_after = 10 if has_intro else 12
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.65),
                Inches(top + 0.08),
                Inches(7.5),
                Inches(avail_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255) # White card
            card.line.color.rgb = RGBColor(226, 232, 240)
            card.line.width = Pt(1.5)
            
            self._add_textbox_column(
                slide,
                body_lines,
                left=Inches(0.95),
                top=Inches(top + 0.2),
                width=Inches(6.9),
                height=Inches(avail_height - 0.25),
                computed_font_size=computed_font_size,
                space_after=space_after
            )

            # Draw illustration card on the right
            self._add_illustration_card(
                slide,
                Inches(8.55),
                Inches(top + 0.08),
                Inches(4.13),
                Inches(avail_height),
                icon_name
            )
        else:
            # Group lines by heading
            groups = []
            current_group = []
            for line in body_lines:
                is_bullet = line.startswith("- ")
                is_heading = (not is_bullet and not line.startswith("  ") and len(line) < 58) or (is_bullet and line.endswith(":") and len(line) < 45)
                
                if is_heading:
                    if current_group:
                        groups.append(current_group)
                    current_group = [line]
                else:
                    if current_group:
                        current_group.append(line)
                    else:
                        current_group = [line]
            if current_group:
                groups.append(current_group)

            # Balanced distribution between two columns
            col1 = []
            col2 = []
            for group in groups:
                if len(col1) <= len(col2):
                    col1.extend(group)
                else:
                    col2.extend(group)

            max_col_lines = max(len(col1), len(col2))
            if max_col_lines <= 5:
                computed_font_size = 15 if has_intro else 16.5
                space_after = 8
            elif max_col_lines <= 8:
                computed_font_size = 13 if has_intro else 14.5
                space_after = 6
            elif max_col_lines <= 12:
                computed_font_size = 11.5 if has_intro else 12.5
                space_after = 5
            elif max_col_lines <= 16:
                computed_font_size = 10.5 if has_intro else 11.5
                space_after = 4
            else:
                computed_font_size = 9.5 if has_intro else 10.5
                space_after = 3

            # Col 1 Card background
            card1 = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.65),
                Inches(top + 0.08),
                Inches(5.8),
                Inches(avail_height)
            )
            card1.fill.solid()
            card1.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card1.line.color.rgb = RGBColor(226, 232, 240)
            card1.line.width = Pt(1.5)
            
            # Col 2 Card background
            if col2:
                card2 = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(6.88),
                    Inches(top + 0.08),
                    Inches(5.8),
                    Inches(avail_height)
                )
                card2.fill.solid()
                card2.fill.fore_color.rgb = RGBColor(255, 255, 255)
                card2.line.color.rgb = RGBColor(226, 232, 240)
                card2.line.width = Pt(1.5)

            self._add_textbox_column(
                slide,
                col1,
                left=Inches(0.95),
                top=Inches(top + 0.25),
                width=Inches(5.2),
                height=Inches(avail_height - 0.3),
                computed_font_size=computed_font_size,
                space_after=space_after
            )
            if col2:
                self._add_textbox_column(
                    slide,
                    col2,
                    left=Inches(7.18),
                    top=Inches(top + 0.25),
                    width=Inches(5.2),
                    height=Inches(avail_height - 0.3),
                    computed_font_size=computed_font_size,
                    space_after=space_after
                )
    def _summarize_raw_result(self, raw_str):
        # Clean text first to see if it looks like table/describe values
        raw_str_clean = raw_str.strip()
        if not raw_str_clean:
            return raw_str
            
        is_raw_table = False
        
        # Check for pandas describe() indicators
        describe_keywords = ["count", "mean", "std", "min", "25%", "50%", "75%", "max"]
        keyword_count = sum(1 for kw in describe_keywords if kw in raw_str_clean.lower())
        if keyword_count >= 3:
            is_raw_table = True
            
        # Check for NaN values or column structure with numbers
        if not is_raw_table:
            lines = raw_str_clean.split("\n")
            if len(lines) >= 3:
                # If a line contains many consecutive spaces and multiple numbers/NaN
                tab_or_space_cols = 0
                for line in lines[:8]:
                    parts = re.split(r'\s{2,}', line.strip())
                    if len(parts) >= 3 and any(p == "NaN" or re.match(r'^-?\d+\.?\d*$', p) for p in parts):
                        tab_or_space_cols += 1
                if tab_or_space_cols >= 2:
                    is_raw_table = True
                    
        if not is_raw_table:
            return raw_str
            
        from core.gemini_service import generate_text
        from core.config import GEMINI_API_KEY
        
        if not GEMINI_API_KEY:
            return raw_str
            
        prompt = (
            f"You are a Senior Business Analyst presenting to corporate stakeholders.\n"
            f"Here is a raw data result of an analysis query:\n\n"
            f"\"\"\"\n{raw_str}\n\"\"\"\n\n"
            f"Please explain and summarize the key business insights and metrics from this raw data "
            f"in 3-4 professional, natural language bullet points (start each with a dash '- '). "
            f"Do NOT output any raw tables, matrices, NaNs, or python structures. "
            f"Convert the raw numbers into meaningful percentage changes, averages, or findings that business stakeholders can easily understand."
        )
        try:
            summary = generate_text(GEMINI_API_KEY, prompt, max_output_tokens=1000)
            if summary and len(summary.strip()) > 10:
                return summary.strip()
        except Exception:
            pass
        return raw_str

    def _generate_dynamic_art(self, seed_string, is_conclusion=False):
        import re
        safe_seed = re.sub(r'[^a-zA-Z0-9_-]', '_', seed_string)
        prefix = "conclusion_" if is_conclusion else "cover_"
        output_path = os.path.join("outputs", f"{prefix}{safe_seed}.png")
        
        if os.path.exists(output_path):
            return output_path
            
        try:
            import math
            import random
            from PIL import Image, ImageDraw, ImageFilter
            
            width, height = 1200, 1200
            image = Image.new("RGBA", (width, height), (15, 23, 42, 255))
            draw = ImageDraw.Draw(image)
            
            random.seed(seed_string)
            
            palettes = [
                ((13, 148, 136), (56, 189, 248), (37, 99, 235), (20, 110, 120)),   # Teal/Cyan/Blue
                ((139, 92, 246), (167, 139, 250), (236, 72, 153), (90, 60, 180)),  # Purple/Indigo/Pink
                ((16, 185, 129), (110, 231, 183), (6, 182, 212), (10, 120, 90)),   # Emerald/Cyan
                ((245, 158, 11), (251, 146, 60), (244, 63, 94), (170, 100, 10))    # Amber/Orange/Rose
            ]
            p_color, s_color, a_color, d_color = random.choice(palettes)
            
            # Background gradient
            for y in range(height):
                for x in range(width):
                    dist = (x + y) / (width + height)
                    r = int(15 + (d_color[0] - 15) * dist * 0.45)
                    g = int(23 + (d_color[1] - 23) * dist * 0.45)
                    b = int(42 + (d_color[2] - 42) * dist * 0.45)
                    image.putpixel((x, y), (r, g, b, 255))
                    
            draw = ImageDraw.Draw(image)
            
            # Faint grid lines
            grid_color = (255, 255, 255, 10)
            vanish_x, vanish_y = width // 2, -height // 4
            for i in range(-10, 21):
                x_target = width * i / 10
                draw.line([(vanish_x, vanish_y), (x_target, height)], fill=grid_color, width=2)
            for i in range(15):
                y = height * (i / 14) ** 1.5
                draw.line([(0, y), (width, y)], fill=grid_color, width=1)
                
            if not is_conclusion:
                # Cover Art
                for w in range(3):
                    points = []
                    amplitude = random.uniform(80, 150)
                    frequency = random.uniform(0.003, 0.006)
                    phase = random.uniform(0, 2 * math.pi)
                    y_offset = height * (0.45 + w * 0.12)
                    
                    for x in range(0, width, 10):
                        y = y_offset + amplitude * math.sin(frequency * x + phase)
                        points.append((x, y))
                        
                    color_w = a_color if w == 0 else (p_color if w == 1 else s_color)
                    draw.line(points, fill=color_w + (30,), width=14)
                    draw.line(points, fill=color_w + (80,), width=6)
                    draw.line(points, fill=(255, 255, 255, 180), width=2)
                    
                # Network Nodes
                nodes = []
                num_nodes = random.randint(12, 18)
                for _ in range(num_nodes):
                    nx = random.uniform(width * 0.15, width * 0.85)
                    ny = random.uniform(height * 0.25, height * 0.8)
                    nr = random.uniform(12, 28)
                    nodes.append((nx, ny, nr))
                    
                for i in range(len(nodes)):
                    dists = []
                    for j in range(len(nodes)):
                        if i != j:
                            d = math.hypot(nodes[i][0]-nodes[j][0], nodes[i][1]-nodes[j][1])
                            dists.append((d, j))
                    dists.sort()
                    for d, idx in dists[:random.randint(1, 2)]:
                        if d < 350:
                            alpha = int(max(10, 140 * (1 - d / 350)))
                            draw.line([(nodes[i][0], nodes[i][1]), (nodes[idx][0], nodes[idx][1])], fill=s_color + (alpha,), width=2)
                            
                for nx, ny, nr in nodes:
                    draw.ellipse([nx - nr*2.2, ny - nr*2.2, nx + nr*2.2, ny + nr*2.2], fill=p_color + (25,))
                    draw.ellipse([nx - nr, ny - nr, nx + nr, ny + nr], fill=s_color + (160,), outline=(255, 255, 255, 200), width=2)
                    draw.ellipse([nx - nr*0.4, ny - nr*0.4, nx + nr*0.4, ny + nr*0.4], fill=(255, 255, 255, 255))
                    
                # Stacked Database Cylinder
                cx, cy, cw, ch = width * 0.5, height * 0.42, 110, 180
                for i in range(4):
                    dy = cy - i * 40
                    draw.ellipse([cx - cw, dy - 25 + 8, cx + cw, dy + 25 + 8], fill=(0,0,0,70))
                    draw.ellipse([cx - cw, dy - 25, cx + cw, dy + 25], fill=p_color + (90,), outline=(255,255,255,160), width=3)
                    draw.ellipse([cx - cw*0.75, dy - 16, cx + cw*0.75, dy + 16], fill=s_color + (140,))
                    
            else:
                # Conclusion Art
                for r in [120, 220, 320, 420]:
                    alpha = int(100 * (1 - r / 500))
                    draw.ellipse([width//2 - r, height//2 - r, width//2 + r, height//2 + r], outline=s_color + (alpha,), width=3)
                    
                num_rays = 12
                for i in range(num_rays):
                    angle = i * (2 * math.pi / num_rays)
                    x1 = width//2 + 90 * math.cos(angle)
                    y1 = height//2 + 90 * math.sin(angle)
                    x2 = width//2 + 380 * math.cos(angle)
                    y2 = height//2 + 380 * math.sin(angle)
                    draw.line([(x1, y1), (x2, y2)], fill=p_color + (40,), width=2)
                    
                cx, cy, cr = width//2, height//2, 85
                draw.ellipse([cx-cr*1.6, cy-cr*1.6, cx+cr*1.6, cy+cr*1.6], fill=a_color + (35,))
                draw.ellipse([cx-cr, cy-cr, cx+cr, cy+cr], fill=p_color + (160,), outline=(255,255,255,210), width=4)
                
                p1 = (cx - cr * 0.4, cy - cr * 0.05)
                p2 = (cx - cr * 0.1, cy + cr * 0.3)
                p3 = (cx + cr * 0.45, cy - cr * 0.35)
                draw.line([p1, p2, p3], fill=(255, 255, 255, 255), width=10, joint="round")
                
            image = image.filter(ImageFilter.SMOOTH_MORE)
            image.save(output_path)
            return output_path
        except Exception:
            return None

    def _add_title_slide(self, prs, file_name):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Dark solid background
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900
        
        # Add the gorgeous cover art image at the right side of the slide
        cover_art_path = self._generate_dynamic_art(file_name, is_conclusion=False)
        if not cover_art_path:
            cover_art_path = self._get_asset_path("cover_illustration.png")
            
        if cover_art_path:
            try:
                slide.shapes.add_picture(
                    cover_art_path,
                    Inches(7.0),
                    Inches(0.5),
                    width=Inches(5.8),
                    height=Inches(6.5)
                )
            except Exception:
                pass
        else:
            # Fallback decorative geometric triangle
            accent_shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                Inches(8.5), Inches(0), Inches(4.833), Inches(7.5)
            )
            accent_shape.fill.solid()
            accent_shape.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
            accent_shape.line.fill.background()
            accent_shape.rotation = 180
        
        # Title text box
        title_box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.2),
            Inches(5.5),
            Inches(2.0),
        )
        frame = title_box.text_frame
        frame.word_wrap = True
        frame.clear()
        p = frame.paragraphs[0]
        p.text = "📊 AI Data Analyst Report"
        p.font.name = "Times New Roman"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Subtitle text box
        subtitle = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(3.5),
            Inches(5.5),
            Inches(1.0),
        )
        frame = subtitle.text_frame
        frame.word_wrap = True
        frame.clear()
        p = frame.paragraphs[0]
        clean_file = os.path.basename(file_name)
        p.text = f"📂 Dataset: {clean_file}"
        p.font.name = "Times New Roman"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
        
        # Descriptive block
        desc_box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(4.8),
            Inches(5.8),
            Inches(1.5),
        )
        frame = desc_box.text_frame
        frame.word_wrap = True
        frame.clear()
        p = frame.paragraphs[0]
        p.text = "✨ Automated profile, deep insights, strategic recommendations, and selected visual evidence."
        p.font.name = "Times New Roman"
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RGBColor(13, 148, 136) # Teal accent

    def _add_dataset_overview(self, prs, profile):
        slide = self._blank_slide(prs, "🔍 Dataset Overview")
        rows = profile.get("rows", 0)
        columns = profile.get("columns", 0)
        duplicates = profile.get("duplicates", 0)
        column_names = profile.get("column_names", [])

        # Format number with commas
        formatted_rows = f"{rows:,}" if isinstance(rows, (int, float)) else str(rows)
        formatted_cols = f"{columns:,}" if isinstance(columns, (int, float)) else str(columns)
        formatted_dups = f"{duplicates:,}" if isinstance(duplicates, (int, float)) else str(duplicates)

        # 3 KPI Cards side-by-side
        self._add_kpi_card(slide, Inches(0.65), Inches(1.4), Inches(3.7), Inches(1.8), "Total Records", formatted_rows, "📊")
        self._add_kpi_card(slide, Inches(4.75), Inches(1.4), Inches(3.7), Inches(1.8), "Total Fields (Columns)", formatted_cols, "🗂️")
        self._add_kpi_card(slide, Inches(8.85), Inches(1.4), Inches(3.7), Inches(1.8), "Duplicate Rows", formatted_dups, "🔄")

        # Column Preview Card below
        col_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.65), Inches(3.55), Inches(7.5), Inches(3.3)
        )
        col_card.fill.solid()
        col_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        col_card.line.color.rgb = RGBColor(226, 232, 240)
        col_card.line.width = Pt(1.5)
        
        preview_box = slide.shapes.add_textbox(
            Inches(0.95), Inches(3.75), Inches(6.9), Inches(2.9)
        )
        tf = preview_box.text_frame
        tf.word_wrap = True
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = "📋 Column Preview"
        p0.font.name = "Times New Roman"
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(30, 41, 59)
        
        p1 = tf.add_paragraph()
        p1.text = "Primary fields detected:"
        p1.font.name = "Times New Roman"
        p1.font.size = Pt(14)
        p1.font.color.rgb = RGBColor(100, 116, 139)
        p1.space_before = Pt(6)
        
        p2 = tf.add_paragraph()
        cols_text = ", ".join(column_names[:16])
        if len(column_names) > 16:
            cols_text += f" (+ {len(column_names) - 16} more)"
        p2.text = cols_text
        p2.font.name = "Times New Roman"
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(71, 85, 105)
        p2.line_spacing = 1.3
        p2.space_before = Pt(8)

        # Visual Illustration Card on the right
        self._add_illustration_card(slide, Inches(8.55), Inches(3.55), Inches(4.13), Inches(3.3), "icon_summary.png")

    def _add_data_quality_slide(self, prs, quality_report):
        slide = self._blank_slide(prs, "🛡️ Data Quality Assessment")
        duplicates = quality_report.get("duplicates", 0)
        constant_columns = len(quality_report.get("constant_columns", []))
        high_cardinality = len(quality_report.get("high_cardinality", {}))
        
        # Safely handle missing values counting
        missing_dict = quality_report.get("missing_values", {})
        if isinstance(missing_dict, dict):
            missing_values = sum(missing_dict.values())
        elif isinstance(missing_dict, (int, float)):
            missing_values = int(missing_dict)
        else:
            missing_values = 0

        # Format metrics
        formatted_dups = f"{duplicates:,}" if isinstance(duplicates, (int, float)) else str(duplicates)
        formatted_miss = f"{missing_values:,}" if isinstance(missing_values, (int, float)) else str(missing_values)
        formatted_const = f"{constant_columns:,}" if isinstance(constant_columns, (int, float)) else str(constant_columns)
        formatted_card = f"{high_cardinality:,}" if isinstance(high_cardinality, (int, float)) else str(high_cardinality)

        # 4 KPI Cards side-by-side
        self._add_kpi_card(slide, Inches(0.65), Inches(1.4), Inches(2.7), Inches(1.8), "Duplicate Rows", formatted_dups, "👥")
        self._add_kpi_card(slide, Inches(3.75), Inches(1.4), Inches(2.7), Inches(1.8), "Missing Cells", formatted_miss, "🔍")
        self._add_kpi_card(slide, Inches(6.85), Inches(1.4), Inches(2.7), Inches(1.8), "Constant Fields", formatted_const, "🔄")
        self._add_kpi_card(slide, Inches(9.95), Inches(1.4), Inches(2.7), Inches(1.8), "High Cardinality", formatted_card, "⚡")

        # Summary box below
        summary_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.65), Inches(3.55), Inches(7.5), Inches(3.3)
        )
        summary_card.fill.solid()
        summary_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        summary_card.line.color.rgb = RGBColor(226, 232, 240)
        summary_card.line.width = Pt(1.5)
        
        summary_box = slide.shapes.add_textbox(
            Inches(0.95), Inches(3.75), Inches(6.9), Inches(2.9)
        )
        tf = summary_box.text_frame
        tf.word_wrap = True
        tf.clear()
        
        p0 = tf.paragraphs[0]
        p0.text = "🛡️ Assessment Summary"
        p0.font.name = "Times New Roman"
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(30, 41, 59)
        
        p1 = tf.add_paragraph()
        p1.text = "• Duplicates: Repeated rows can lead to statistical bias and should be cleaned."
        p1.font.name = "Times New Roman"
        p1.font.size = Pt(13.5)
        p1.font.color.rgb = RGBColor(71, 85, 105)
        p1.space_before = Pt(6)
        
        p2 = tf.add_paragraph()
        p2.text = "• Missing Values: High concentrations of null data require appropriate imputation or removal."
        p2.font.name = "Times New Roman"
        p2.font.size = Pt(13.5)
        p2.font.color.rgb = RGBColor(71, 85, 105)
        p2.space_before = Pt(4)
        
        p3 = tf.add_paragraph()
        p3.text = "• Constant & Cardinality: Single-value columns contain no descriptive power; extremely high unique columns should be modeled carefully."
        p3.font.name = "Times New Roman"
        p3.font.size = Pt(13.5)
        p3.font.color.rgb = RGBColor(71, 85, 105)
        p3.space_before = Pt(4)

        # Visual Illustration Card on the right
        self._add_illustration_card(slide, Inches(8.55), Inches(3.55), Inches(4.13), Inches(3.3), "icon_risks.png")

    def _add_analysis_slide(self, prs, analysis_result):
        if isinstance(analysis_result, dict):
            result_text = analysis_result.get("result", "")
        else:
            result_text = analysis_result

        if not self._is_valid_content(result_text):
            return

        # Clean/Format raw result
        if isinstance(result_text, pd.DataFrame):
            raw_str = result_text.to_string(index=True)
        elif isinstance(result_text, pd.Series):
            raw_str = result_text.to_string()
        else:
            raw_str = str(result_text)

        # Call the summarizer helper to turn raw table data into readable business bullets
        final_text = self._summarize_raw_result(raw_str)

        base_name = getattr(self, "_current_base_name", "")
        analysis_default = f"{base_name} Analysis" if base_name else "Analysis Results"
        analysis_title = self._determine_slide_title(final_text, analysis_default)

        # Make sure title is styled with emoji
        final_title = analysis_title if analysis_title.startswith("📈") else f"📈 {analysis_title}"
        self._add_text_slides(prs, final_title, final_text)

    def _add_visualization_slides(self, prs, chart_items=None):
        if not chart_items:
            return

        for i, item in enumerate(chart_items[:6], start=1):
            try:
                path, title, caption = item
            except Exception:
                try:
                    path, caption = item
                    title = self._title_from_caption(caption, f"Visualization {i}")
                except Exception:
                    path = item
                    title = f"Visualization {i}"
                    caption = ""

            if not path or not os.path.exists(path):
                continue

            # Ensure header has visualization emoji
            slide_title = title if title.startswith("🖼️") else f"🖼️ {title}"
            slide = self._blank_slide(prs, slide_title)
            
            # Draw a beautiful card shape behind the chart
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5),
                Inches(1.35),
                Inches(10.33),
                Inches(5.45)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            card.line.width = Pt(1.5)

            try:
                slide.shapes.add_picture(
                    path,
                    Inches(2.16),
                    Inches(1.55),
                    width=Inches(9.0),
                    height=Inches(4.15),
                )
            except Exception:
                continue

            cap = slide.shapes.add_textbox(
                Inches(1.8),
                Inches(5.8),
                Inches(9.73),
                Inches(0.8),
            )
            frame = cap.text_frame
            frame.clear()
            frame.word_wrap = True
            p = frame.paragraphs[0]
            display_caption = self._clean_text(caption)[0] if self._clean_text(caption) else str(title)
            display_caption = re.sub(r"\*\*(.*?)\*\*", r"\1", display_caption)
            display_caption = re.sub(r"^#+\s*", "", display_caption)
            p.text = display_caption
            p.font.name = "Times New Roman"
            p.font.size = Pt(14.5)
            p.font.color.rgb = RGBColor(71, 85, 105)
            p.alignment = PP_ALIGN.CENTER

    def _add_text_slides(self, prs, title, text):
        if not self._is_valid_content(text):
            return

        # Clean/Format raw result
        if isinstance(text, pd.DataFrame):
            raw_str = text.to_string(index=True)
        elif isinstance(text, pd.Series):
            raw_str = text.to_string()
        else:
            raw_str = str(text)

        # Summarize raw table/describe numerical values automatically
        text = self._summarize_raw_result(raw_str)

        lines = self._clean_text(text)
        if not lines:
            return

        sections = []
        current_title = title
        current_lines = []

        heading_patterns = [
            r'^#{1,4}\s+(.+)$',
            r'^\*\*(.+?)\*\*$',
            r'^\*\*(.+?):\*\*$'
        ]

        def extract_heading(line):
            for pattern in heading_patterns:
                match = re.match(pattern, line)
                if match:
                    heading = match.group(1).strip()
                    if len(heading) < 80:
                        return heading
            return None

        for line in lines:
            heading = extract_heading(line)
            if heading:
                if current_lines:
                    sections.append((current_title, current_lines))
                current_title = heading
                current_lines = []
            else:
                # avoid duplicating the section title in the very first line of a section
                clean_line = re.sub(r'^#+\s*', '', line).strip(":- ")
                clean_line = re.sub(r"\*\*(.*?)\*\*", r"\1", clean_line)
                if not current_lines and (
                    clean_line.lower() == current_title.lower() or
                    current_title.lower() in clean_line.lower() or
                    clean_line.lower() in current_title.lower()
                ):
                    continue
                current_lines.append(line)

        if current_lines:
            sections.append((current_title, current_lines))

        for section_title, section_lines in sections:
            if not section_lines:
                continue

            chunks = self._split_section_into_chunks(section_lines, max_chars=800)

            for idx, chunk in enumerate(chunks):
                # Ensure the section title has a clean emoji representing it
                decorated_title = section_title
                if not any(char in decorated_title for char in ["📊", "🔍", "🛡️", "📈", "💡", "🎯", "🏁", "🖼️"]):
                    emoji_prefix = self._get_emoji_for_text(decorated_title)
                    if emoji_prefix and emoji_prefix != "•":
                        decorated_title = f"{emoji_prefix} {decorated_title}"
                    else:
                        decorated_title = f"📄 {decorated_title}"

                slide_title = (
                    decorated_title
                    if idx == 0
                    else f"{decorated_title} (Continued)"
                )
                slide = self._blank_slide(prs, slide_title)
                self._add_body_lines(slide, chunk, top=1.22)

    def _add_conclusion_slide(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900
        
        # Decorative shape
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE,
            Inches(0), Inches(0), Inches(4.5), Inches(7.5)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
        accent_shape.line.fill.background()
        
        # Center the conclusion illustration inside the left dark panel
        base_name = getattr(self, "_current_base_name", "")
        conclusion_icon_path = self._generate_dynamic_art(base_name or "conclusion", is_conclusion=True)
        if not conclusion_icon_path:
            conclusion_icon_path = self._get_asset_path("icon_conclusion.png")
            
        if conclusion_icon_path:
            try:
                slide.shapes.add_picture(
                    conclusion_icon_path,
                    Inches(0.5),
                    Inches(2.0),
                    width=Inches(3.5)
                )
            except Exception:
                pass
        
        # Text box on the right
        text_box = slide.shapes.add_textbox(
            Inches(5.0),
            Inches(2.2),
            Inches(7.5),
            Inches(4.0),
        )
        frame = text_box.text_frame
        frame.word_wrap = True
        frame.clear()
        
        p = frame.paragraphs[0]
        p.text = "🏁 Conclusion & Next Steps"
        p.font.name = "Times New Roman"
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        p2 = frame.add_paragraph()
        p2.text = "This report was automatically generated by the AI Data Analyst Agent."
        p2.font.name = "Times New Roman"
        p2.font.size = Pt(19)
        p2.font.color.rgb = RGBColor(148, 163, 184)
        p2.space_before = Pt(20)
        
        p3 = frame.add_paragraph()
        p3.text = "🔄 The workflow included data quality assessment, dataset profiling, AI-powered analysis, visualization, business insights, and strategic recommendations."
        p3.font.name = "Times New Roman"
        p3.font.size = Pt(19)
        p3.font.color.rgb = RGBColor(148, 163, 184)
        p3.space_before = Pt(12)
        
        p4 = frame.add_paragraph()
        p4.text = "🚀 Use these findings to drive strategic business decisions and optimize your processes."
        p4.font.name = "Times New Roman"
        p4.font.size = Pt(19)
        p4.font.bold = True
        p4.font.color.rgb = RGBColor(13, 148, 136) # Teal
        p4.space_before = Pt(16)

    def _split_section_into_chunks(self, lines, max_chars=800):
        """
        Split a section into multiple slides only if it becomes too large.
        Keeps headings together.
        """
        chunks = []
        current_chunk = []
        current_size = 0

        for line in lines:
            line_size = len(line)

            if current_size + line_size > max_chars and current_chunk:
                chunks.append(current_chunk)
                current_chunk = [line]
                current_size = line_size
            else:
                current_chunk.append(line)
                current_size += line_size

        if current_chunk:
            chunks.append(current_chunk)

        return chunks

    def _remove_empty_slides(self, prs):
        slides_to_delete = []
        for idx, slide in enumerate(prs.slides):
            text_found = False
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if len(txt) > 5:
                        text_found = True
                        break
            if not text_found:
                slides_to_delete.append(idx)

        for idx in reversed(slides_to_delete):
            slide_id = prs.slides._sldIdLst[idx]
            rel_id = slide_id.rId
            prs.part.drop_rel(rel_id)
            del prs.slides._sldIdLst[idx]
